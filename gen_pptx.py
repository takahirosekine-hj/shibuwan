#!/usr/bin/env python3
"""Generate 渋一プロジェクト概要.pptx - 16:9 Dark Theme Presentation"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# Constants
W, H = Inches(16), Inches(9)
BG = RGBColor(0x0E, 0x0E, 0x12)
RED = RGBColor(0xCC, 0x00, 0x00)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
DARK = RGBColor(0x12, 0x12, 0x16)
GRAY = RGBColor(0x88, 0x88, 0x88)
GOLD = RGBColor(0xD4, 0xA8, 0x45)
LIGHT_GRAY = RGBColor(0xAA, 0xAA, 0xAA)
CARD_BG = RGBColor(0x1A, 0x1A, 0x22)
FONT = 'Hiragino Kaku Gothic ProN'

prs = Presentation()
prs.slide_width = W
prs.slide_height = H


def set_bg(slide, color=BG):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_shape(slide, left, top, width, height, fill_color=None, line_color=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.line.fill.background()
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = Pt(1)
    return shape


def add_text(slide, left, top, width, height, text, size=18, color=WHITE, bold=False, align=PP_ALIGN.LEFT, font_name=FONT):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = align
    return txBox


def add_multi_text(slide, left, top, width, height, lines):
    """lines: list of (text, size, color, bold, align)"""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        text, size, color, bold, align = line
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = text
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.font.bold = bold
        p.font.name = FONT
        p.alignment = align
    return txBox


# ─── SLIDE 1: TITLE ───
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_shape(slide, Inches(0), Inches(0), W, Inches(0.08), fill_color=RED)
add_multi_text(slide, Inches(1.5), Inches(1.8), Inches(13), Inches(6), [
    ("渋一（シブワン）", 52, WHITE, True, PP_ALIGN.CENTER),
    ("", 12, WHITE, False, PP_ALIGN.CENTER),
    ("渋谷発メンズIP創出プロジェクト", 28, GRAY, False, PP_ALIGN.CENTER),
    ("Shibuya Ikemen No.1 Project  \u00d7  ストーリーコンテンツ", 20, LIGHT_GRAY, False, PP_ALIGN.CENTER),
    ("", 16, WHITE, False, PP_ALIGN.CENTER),
    ("企画書  プロジェクト概要", 24, GOLD, True, PP_ALIGN.CENTER),
])
add_shape(slide, Inches(5), Inches(6.2), Inches(6), Inches(0.04), fill_color=RED)
add_text(slide, Inches(1.5), Inches(7.5), Inches(13), Inches(0.6),
         "CONFIDENTIAL  \u00a9  HJ Inc. All Rights Reserved.", 14, GRAY, align=PP_ALIGN.CENTER)

# ─── SLIDE 2: MISSION ───
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_shape(slide, Inches(0), Inches(0), Inches(0.12), H, fill_color=RED)
add_text(slide, Inches(0.8), Inches(0.5), Inches(5), Inches(0.6), "MISSION", 14, RED, bold=True)
add_text(slide, Inches(0.8), Inches(1.0), Inches(14), Inches(1.5),
         "Z\u30fb\u03b1\u4e16\u4ee3\u5973\u6027\u304c\u300c\u6700\u521d\u306b\u63a8\u3057\u305f\u7537\u6027\u30b9\u30bf\u30fc\u300d\u3092\u6bce\u5e74\u751f\u307f\u51fa\u3057\u7d9a\u3051\u308b\u3001\n\u65e5\u672c\u6700\u5927\u306e\u30e1\u30f3\u30baIP\u5275\u51fa\u88c5\u7f6e\u3002",
         22, WHITE, bold=True)

pillars = [
    ("No.1 \u30dd\u30b8\u30b7\u30e7\u30f3\u78ba\u7acb", "\u30e1\u30f3\u30ba\u30a2\u30a4\u30c9\u30eb\u5e02\u5834\u306b\u304a\u3051\u308b\n\u5727\u5012\u7684\u30b7\u30a7\u30a2\u3092\u7372\u5f97", RED),
    ("\u5727\u5012\u7684\u306a\u8a71\u984c\u6027", "SNS\u30d0\u30ba\u3001\u30e1\u30c7\u30a3\u30a2\u9732\u51fa\u3001\n\u30c8\u30ec\u30f3\u30c9\u5165\u308a\u3092\u5e38\u614b\u5316", GOLD),
    ("\u6700\u5f37\u306e\u30d5\u30a1\u30f3\u30c0\u30e0", "\u71b1\u91cf\u306e\u9ad8\u3044\u30d5\u30a1\u30f3\u30b3\u30df\u30e5\u30cb\u30c6\u30a3\u3068\n\u7d99\u7d9a\u7684\u306a\u7d4c\u6e08\u570f\u3092\u5275\u51fa", RGBColor(0x1E, 0x6F, 0xBF)),
]
for i, (title, desc, accent) in enumerate(pillars):
    x = Inches(0.8 + i * 4.9)
    add_shape(slide, x, Inches(3.2), Inches(4.4), Inches(3.0), fill_color=CARD_BG)
    add_shape(slide, x, Inches(3.2), Inches(4.4), Inches(0.06), fill_color=accent)
    add_text(slide, x + Inches(0.3), Inches(3.5), Inches(3.8), Inches(0.6), title, 20, accent, bold=True)
    add_text(slide, x + Inches(0.3), Inches(4.3), Inches(3.8), Inches(1.5), desc, 16, LIGHT_GRAY)

add_shape(slide, Inches(0.8), Inches(7.0), Inches(14.4), Inches(0.8), fill_color=RGBColor(0x1A, 0x0A, 0x0A))
add_text(slide, Inches(1.2), Inches(7.1), Inches(13.6), Inches(0.6),
         "\u76ee\u6a19:  2028\u5e74  \u30e1\u30f3\u30ba\u82b8\u80fd\u90e8\u9580\uff11\u4f4d  \uff1d  100\u5104\u5186IP\u306e\u5b9f\u73fe", 20, RED, bold=True, align=PP_ALIGN.CENTER)

# ─── SLIDE 3: CONCEPT ───
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_shape(slide, Inches(0), Inches(0), Inches(0.12), H, fill_color=RED)
add_text(slide, Inches(0.8), Inches(0.5), Inches(5), Inches(0.6), "CONCEPT", 14, RED, bold=True)
add_multi_text(slide, Inches(0.8), Inches(1.0), Inches(14), Inches(2), [
    ("\u300c\u653e\u8ab2\u5f8c\u6f14\u6280\u6d3e\u30af\u30e9\u30d6\u300d\u2500\u2500 \u30b9\u30c8\u30fc\u30ea\u30fc \u00d7 \u30a2\u30a4\u30c9\u30eb", 26, WHITE, True, PP_ALIGN.LEFT),
    ("", 8, WHITE, False, PP_ALIGN.LEFT),
    ("\u6e0b\u8c37\u306e10\u306e\u67b6\u7a7a\u9ad8\u6821\u3092\u821e\u53f0\u306b\u3001\u5404\u6821\u304c\u30e1\u30f3\u30ba\u30dc\u30fc\u30ab\u30eb\u30e6\u30cb\u30c3\u30c8\u3092\u7d50\u6210\u3002", 18, LIGHT_GRAY, False, PP_ALIGN.LEFT),
    ("\u30b7\u30e7\u30fc\u30c8\u30c9\u30e9\u30de\u3067\u63cf\u304f\u9752\u6625\u30b9\u30c8\u30fc\u30ea\u30fc\u3068\u3001\u30ea\u30a2\u30eb\u306e\u30d1\u30d5\u30a9\u30fc\u30de\u30f3\u30b9\u304c\u4ea4\u5dee\u3059\u308b\u65b0\u4e16\u4ee3\u30a8\u30f3\u30bf\u30e1\u3002", 18, LIGHT_GRAY, False, PP_ALIGN.LEFT),
])

axes = [
    ("\u89aa\u8fd1\u611f\u8ef8\uff08\u30b3\u30e0\u30c9\u30c3\u30c8\u7684\uff09", [
        "\u65e5\u5e38\u30fb\u30d0\u30e9\u30a8\u30c6\u30a3\u4f01\u753b",
        "\u30aa\u30d5\u30b7\u30e7\u30c3\u30c8\u30fb\u7d20\u306e\u8868\u60c5",
        "\u30e1\u30f3\u30d0\u30fc\u9593\u306e\u95a2\u4fc2\u6027",
        "\u5207\u308a\u629c\u304d\u62e1\u6563\u30fb\u30b7\u30e7\u30fc\u30c8\u52d5\u753b",
    ], GOLD),
    ("\u30a2\u30a4\u30c9\u30eb\u8ef8\uff08\u30b9\u30bf\u30fc\u6027\u78ba\u7acb\uff09", [
        "\u697d\u66f2\u30fb\u30c0\u30f3\u30b9\u30d1\u30d5\u30a9\u30fc\u30de\u30f3\u30b9",
        "\u30c9\u30e9\u30de\u30fb\u30b9\u30c8\u30fc\u30ea\u30fc\u3078\u306e\u611f\u60c5\u79fb\u5165",
        "\u30d3\u30b8\u30e5\u30a2\u30eb\u91cd\u8996\u306e\u9ad8\u54c1\u8cea\u30b3\u30f3\u30c6\u30f3\u30c4",
        "\u30d5\u30a1\u30f3\u4ea4\u6d41\u30a4\u30d9\u30f3\u30c8\u30fb\u30e9\u30a4\u30d6",
    ], RED),
]
for i, (title, items, accent) in enumerate(axes):
    x = Inches(0.8 + i * 7.4)
    add_shape(slide, x, Inches(4.0), Inches(6.8), Inches(4.2), fill_color=CARD_BG)
    add_shape(slide, x, Inches(4.0), Inches(6.8), Inches(0.06), fill_color=accent)
    add_text(slide, x + Inches(0.4), Inches(4.3), Inches(6), Inches(0.6), title, 20, accent, bold=True)
    for j, item in enumerate(items):
        add_text(slide, x + Inches(0.6), Inches(5.1 + j * 0.65), Inches(5.8), Inches(0.5),
                 "\u25cf  " + item, 16, LIGHT_GRAY)

# ─── SLIDE 4: WORLD BUILDING ───
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_shape(slide, Inches(0), Inches(0), Inches(0.12), H, fill_color=RED)
add_text(slide, Inches(0.8), Inches(0.5), Inches(8), Inches(0.6), "WORLD BUILDING", 14, RED, bold=True)
add_text(slide, Inches(0.8), Inches(1.0), Inches(14), Inches(0.8),
         "\u6e0b\u8c3710\u6821  \u2500\u2500  10\u306e\u5b66\u5712\u300110\u306e\u30ab\u30eb\u30c1\u30e3\u30fc\u300110\u306e\u30ea\u30fc\u30c0\u30fc", 24, WHITE, bold=True)

schools = [
    ("\u6e0b\u4e00", "\u5b87\u7530\u5ddd\u753a", "\u30b0\u30e9\u30d5\u30a3\u30c6\u30a3", "#CC0000", "\u308b\u3044\uff08\u4e3b\u4eba\u516c\uff09"),
    ("\u767e\u5b9f", "\u767e\u8ed2\u5e97", "\u30a2\u30f3\u30b0\u30e9", "#333333", "\u3086\u3046\u3059\u3051"),
    ("\u30d5\u30a1\u30a4\u30e4\u30fc", "\u30d5\u30a1\u30a4\u30e4\u30fc\u901a\u308a", "\u4e00\u5339\u72fc\u30fb\u6700\u5f37", "#6B2D8B", "\u3044\u304a\u3046"),
    ("\u30df\u30e4\u30b7\u30bf", "\u5bae\u4e0b\u516c\u5712", "\u30b9\u30b1\u30fc\u30bf\u30fc\u30fb\u81ea\u7531", "#1E6FBF", "\u30ab\u30a4"),
    ("\u30aa\u30af\u30b7\u30d6", "\u5965\u6e0b\u8c37", "\u30a8\u30ea\u30fc\u30c8\u30fb\u7ba1\u7406", "#1A5C30", "\u3072\u304b\u308b"),
    ("\u30b5\u30af\u30e9", "\u685c\u30f6\u4e18", "\u9ad8\u8cb4\u30fb\u5de8\u5927\u8cc7\u672c", "#8B6914", "\u30c4\u30e8\u30b7\uff06\u30d2\u30c8\u30df"),
    ("\u30de\u30eb\u30e4\u30de", "\u5186\u5c71\u753a", "\u30af\u30e9\u30d6\u30fbDJ", "#D4267E", "\u30ca\u30aa"),
    ("\u30c9\u30a6\u30b2\u30f3", "\u9053\u7384\u5742", "\u30ab\u30aa\u30b9\u30fb\u96d1\u591a", "#D4742D", "\u30c6\u30c4"),
    ("\u30b3\u30a6\u30a8\u30f3", "\u516c\u5712\u901a\u308a", "\u30a2\u30fc\u30c8\u30fb\u6f14\u5287", "#888888", "\u30ec\u30f3"),
    ("\u30df\u30e4\u30de\u30b9", "\u5bae\u76ca\u5742", "\u30c6\u30c3\u30af\u30fb\u30c7\u30fc\u30bf", "#C4A800", "\u30b7\u30e5\u30f3"),
]

for idx, (name, area, culture, hex_color, leader) in enumerate(schools):
    col = idx % 5
    row = idx // 5
    x = Inches(0.8 + col * 3.0)
    y = Inches(2.2 + row * 3.2)
    r, g, b = int(hex_color[1:3], 16), int(hex_color[3:5], 16), int(hex_color[5:7], 16)
    accent = RGBColor(r, g, b)
    add_shape(slide, x, y, Inches(2.7), Inches(2.8), fill_color=CARD_BG)
    add_shape(slide, x, y, Inches(2.7), Inches(0.06), fill_color=accent)
    add_text(slide, x + Inches(0.2), y + Inches(0.2), Inches(2.3), Inches(0.5), name, 20, accent, bold=True)
    add_text(slide, x + Inches(0.2), y + Inches(0.75), Inches(2.3), Inches(0.4), area, 13, LIGHT_GRAY)
    add_text(slide, x + Inches(0.2), y + Inches(1.2), Inches(2.3), Inches(0.4), culture, 14, WHITE)
    add_text(slide, x + Inches(0.2), y + Inches(1.75), Inches(2.3), Inches(0.5), leader, 13, GRAY)

# ─── SLIDE 5: FIRST SEASON STORY ───
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_shape(slide, Inches(0), Inches(0), Inches(0.12), H, fill_color=RED)
add_text(slide, Inches(0.8), Inches(0.5), Inches(8), Inches(0.6), "FIRST SEASON STORY", 14, RED, bold=True)
add_multi_text(slide, Inches(0.8), Inches(1.0), Inches(14), Inches(1.5), [
    ("\u30d5\u30a1\u30fc\u30b9\u30c8\u30b7\u30fc\u30ba\u30f3\u516820\u8a71\u300c\u59cb\u307e\u308a\u306e\u708e\u300d", 24, WHITE, True, PP_ALIGN.LEFT),
    ("", 6, WHITE, False, PP_ALIGN.LEFT),
    ("\u5e8f\u7ae0 \u2192 \u30ea\u30f3\u306e\u8ee2\u6821 \u2192 \u6e0b\u4e00vs\u767e\u5b9f\u306e\u5bfe\u7acb\u6fc0\u5316 \u2192 \u5185\u90e8\u5d29\u58ca\u3068\u518d\u8d77 \u2192 \u6700\u7d42\u6c7a\u6226\u3068\u65b0\u305f\u306a\u8105\u5a01", 16, GRAY, False, PP_ALIGN.LEFT),
])

acts = [
    ("\u7b2c\u4e00\u5e55\uff081\u301c5\u8a71\uff09", "\u5747\u8861\u3068\u51fa\u4f1a\u3044", "\u308b\u3044\u3068\u30ea\u30f3\u306e\u51fa\u4f1a\u3044\u3002\u767e\u5b9f\u306e\u3086\u3046\u3059\u3051\u304c\u30ea\u30f3\u306e\u596a\u9084\u3092\u5ba3\u8a00\u3057\u3001\u51b7\u6226\u304c\u5d29\u308c\u308b\u3002", RGBColor(0x1E, 0x6F, 0xBF)),
    ("\u7b2c\u4e8c\u5e55\uff086\u301c10\u8a71\uff09", "\u6fc0\u7a81\u3068\u62e1\u5927", "\u30ad\u30e3\u30c3\u30c8\u30b9\u30c8\u30ea\u30fc\u30c8\u3067\u6e0b\u4e00vs\u767e\u5b9f\u304c\u6fc0\u7a81\u3002\u3044\u304a\u3046\u53c2\u6226\u3002\u3086\u3046\u305f\u306e\u88cf\u5207\u308a\u306e\u7a2e\u304c\u82bd\u751f\u3048\u308b\u3002", GOLD),
    ("\u7b2c\u4e09\u5e55\uff0811\u301c15\u8a71\uff09", "\u5d29\u58ca\u3068\u7d76\u671b", "\u3086\u3046\u305f\u306e\u88cf\u5207\u308a\u3067\u6e0b\u4e00\u5d29\u58ca\u3002Mr.K\u3068\u306e\u9082\u9005\u3067\u308b\u3044\u304c\u518d\u8d77\u3002O.K.thanks\u65b0\u66f2\u304c\u6d41\u308c\u308b\u3002", RED),
    ("\u7b2c\u56db\u5e55\uff0816\u301c20\u8a71\uff09", "\u6700\u7d42\u6c7a\u6226\u3068\u672a\u6765", "\u5168\u9762\u6226\u4e89\u3002\u308b\u3044vs\u3086\u3046\u3059\u3051\u6700\u7d42\u6c7a\u6226\u3067\u52dd\u5229\u3002\u30b5\u30af\u30e9\u304c\u6b21\u306a\u308b\u8105\u5a01\u3068\u3057\u3066\u767b\u5834\u3002", RGBColor(0x1A, 0x5C, 0x30)),
]

for i, (ep, title, desc, accent) in enumerate(acts):
    x = Inches(0.8 + i * 3.7)
    y = Inches(3.2)
    add_shape(slide, x, y, Inches(3.4), Inches(4.8), fill_color=CARD_BG)
    add_shape(slide, x, y, Inches(3.4), Inches(0.06), fill_color=accent)
    add_text(slide, x + Inches(0.25), y + Inches(0.3), Inches(2.9), Inches(0.4), ep, 14, accent, bold=True)
    add_text(slide, x + Inches(0.25), y + Inches(0.8), Inches(2.9), Inches(0.5), title, 22, WHITE, bold=True)
    add_text(slide, x + Inches(0.25), y + Inches(1.6), Inches(2.9), Inches(2.8), desc, 15, LIGHT_GRAY)

# ─── SLIDE 6: CAST ───
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_shape(slide, Inches(0), Inches(0), Inches(0.12), H, fill_color=RED)
add_text(slide, Inches(0.8), Inches(0.5), Inches(5), Inches(0.6), "CAST", 14, RED, bold=True)
add_text(slide, Inches(0.8), Inches(1.0), Inches(14), Inches(0.8),
         "\u30ad\u30e3\u30b9\u30c8  \u2500\u2500  11\u540d + \u5973\u6027\u30ad\u30e3\u30b9\u30c8", 24, WHITE, bold=True)

cast_data = [
    ("\u308b\u3044", "\u6e0b\u4e00", "\u4e3b\u4eba\u516c\u30fb\u4e0d\u5c48\u306e\u7cbe\u795e", RED),
    ("\u306f\u308b", "\u6e0b\u4e00", "\u76f8\u68d2\u30fb\u7279\u653b\u968a\u9577", RED),
    ("\u308a\u3087\u3046\u3059\u3051", "\u6e0b\u4e00", "\u5be1\u9ed9\u306a\u6700\u5f37\u30a4\u30b1\u30e1\u30f3", RED),
    ("\u3086\u3046\u305f", "\u6e0b\u4e00", "\u53c2\u8b00\u2192\u88cf\u5207\u308a", RED),
    ("\u3057\u304a\u3093", "\u6e0b\u4e00", "\u5144\u8cb4\u808c\u2192\u3086\u3046\u305f\u306b\u540c\u8abf", RED),
    ("\u3086\u3046\u3059\u3051", "\u767e\u5b9f", "\u30ea\u30fc\u30c0\u30fc\u30fb\u30ab\u30ea\u30b9\u30de", GRAY),
    ("\u3068\u3046\u307e", "\u767e\u5b9f", "\u526f\u30ea\u30fc\u30c0\u30fc", GRAY),
    ("\u306f\u308c", "\u767e\u5b9f", "\u5207\u308a\u8fbc\u307f\u968a\u9577", GRAY),
    ("\u3053\u3046\u306e\u3059\u3051", "\u767e\u5b9f", "\u60c5\u5831\u5c4b", GRAY),
    ("\u305b\u308a\u306a", "\u767e\u5b9f", "\u30ae\u30e3\u30eb\u30fb\u7d05\u4e00\u70b9", GRAY),
    ("\u3044\u304a\u3046", "\u30d5\u30a1\u30a4\u30e4\u30fc", "\u4e00\u5339\u72fc\u30fb\u6e0b\u8c37\u6700\u5f37", RGBColor(0x6B, 0x2D, 0x8B)),
    ("\u3072\u304b\u308b", "\u30aa\u30af\u30b7\u30d6", "\u7ba1\u7406\u8005\u30ea\u30fc\u30c0\u30fc", RGBColor(0x1A, 0x5C, 0x30)),
]

for i, (name, school, role, accent) in enumerate(cast_data):
    col = i % 4
    row = i // 4
    x = Inches(0.8 + col * 3.7)
    y = Inches(2.2 + row * 2.0)
    add_shape(slide, x, y, Inches(3.4), Inches(1.6), fill_color=CARD_BG)
    add_shape(slide, x, y, Inches(0.08), Inches(1.6), fill_color=accent)
    add_text(slide, x + Inches(0.3), y + Inches(0.15), Inches(3.0), Inches(0.5), name, 20, WHITE, bold=True)
    add_text(slide, x + Inches(0.3), y + Inches(0.65), Inches(3.0), Inches(0.4), school, 13, accent)
    add_text(slide, x + Inches(0.3), y + Inches(1.05), Inches(3.0), Inches(0.4), role, 14, LIGHT_GRAY)

add_text(slide, Inches(0.8), Inches(8.0), Inches(14), Inches(0.5),
         "Female cast (tentative):  \u308a\u3042\u3001\u306f\u308b\u306a\u3001\u307e\u3046\u3001\u306f\u308b\u3042\u3001\u3058\u3085\u306a\u3001\u306a\u308b\u305b", 14, GRAY)

# ─── SLIDE 7: O.K.thanks ───
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_shape(slide, Inches(0), Inches(0), Inches(0.12), H, fill_color=GOLD)
add_text(slide, Inches(0.8), Inches(0.5), Inches(8), Inches(0.6), "O.K.thanks", 14, GOLD, bold=True)
add_multi_text(slide, Inches(0.8), Inches(1.0), Inches(14), Inches(2.5), [
    ("\u4f1d\u8aac\u306e\u5b58\u5728  O.K.thanks", 26, WHITE, True, PP_ALIGN.LEFT),
    ("", 8, WHITE, False, PP_ALIGN.LEFT),
    ("\u304b\u3064\u3066\u6e0b\u8c37\u3092\u7d71\u4e00\u3057\u305f\u30c1\u30fc\u30e0\u3068\u5bfe\u7b49\u306b\u6e21\u308a\u5408\u3063\u305f\u4f1d\u8aac\u7684\u5b58\u5728\u3002", 18, LIGHT_GRAY, False, PP_ALIGN.LEFT),
    ("\u7269\u8a9e\u306e\u4e2d\u3067\u306f\u300c\u30b9\u30da\u30b7\u30e3\u30eb\u30b5\u30f3\u30af\u30b9\u300d\u7684\u306b\u767b\u5834\u3057\u3001\u4e3b\u4eba\u516c\u308b\u3044\u306e\u5fc3\u306b\u706b\u3092\u706f\u3059\u3002", 18, LIGHT_GRAY, False, PP_ALIGN.LEFT),
    ("\u30ea\u30a2\u30eb\u3067\u306f\u30d7\u30ed\u30b8\u30a7\u30af\u30c8\u306e\u30b9\u30dd\u30f3\u30b5\u30fc\u3068\u3057\u3066\u3001\u65b0\u66f2\u304c\u30d5\u30a1\u30fc\u30b9\u30c8\u30b7\u30fc\u30ba\u30f3\u306e\u30c6\u30fc\u30de\u66f2\u306b\u3002", 18, LIGHT_GRAY, False, PP_ALIGN.LEFT),
])

chars = [
    ("Mr.K\uff08\u706b\u3092\u706f\u3059\u8005\uff09", "\u7b2c14\u8a71\u3067\u308b\u3044\u304c\u7d76\u671b\u306e\u5e95\u306b\u3044\u308b\u3068\u304d\u3001\u8def\u5730\u88cf\u3067\u51fa\u4f1a\u3046\u3002",
     "\u300c\u6e0b\u8c37\u306f\u306a\u3001\u8ab0\u304b\u304c\u5b88\u308d\u3046\u3068\u3059\u308b\u9650\u308a\u3001\u6b7b\u306a\u306d\u3047\u3088\u300d", RED),
    ("Mr.O\uff08\u98a8\u3092\u9001\u308b\u8005\uff09", "\u7b2c15\u8a71\u3067\u308b\u3044\u306e\u518d\u8d77\u3092\u5f8c\u62bc\u3057\u3002",
     "\u300c\u3044\u3044\u9854\u3057\u3066\u3093\u306a\u300d", GOLD),
]
for i, (name, scene, quote, accent) in enumerate(chars):
    x = Inches(0.8 + i * 7.4)
    y = Inches(4.5)
    add_shape(slide, x, y, Inches(6.8), Inches(3.5), fill_color=CARD_BG)
    add_shape(slide, x, y, Inches(6.8), Inches(0.06), fill_color=accent)
    add_text(slide, x + Inches(0.4), y + Inches(0.3), Inches(6), Inches(0.5), name, 22, accent, bold=True)
    add_text(slide, x + Inches(0.4), y + Inches(1.0), Inches(6), Inches(0.8), scene, 16, LIGHT_GRAY)
    add_text(slide, x + Inches(0.4), y + Inches(2.0), Inches(6), Inches(1.0), quote, 20, WHITE, bold=True)

# ─── SLIDE 8: ROADMAP ───
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_shape(slide, Inches(0), Inches(0), Inches(0.12), H, fill_color=RED)
add_text(slide, Inches(0.8), Inches(0.5), Inches(5), Inches(0.6), "ROADMAP", 14, RED, bold=True)
add_text(slide, Inches(0.8), Inches(1.0), Inches(14), Inches(0.8),
         "3\u30d5\u30a7\u30fc\u30ba  \u2500\u2500  SNS\u59cb\u52d5\u304b\u3089No.1\u3078", 24, WHITE, bold=True)

phases = [
    ("PHASE 1", "2026", "SNS\u59cb\u52d5\u30fb\u8a8d\u77e5\u62e1\u5927",
     "\u300c\u5b66\u5712\u300d\u30b7\u30e7\u30fc\u30c8\u30e0\u30fc\u30d3\u30fc\u3092\u8ef8\u306b\u30ad\u30e3\u30e9\u30af\u30bf\u30fc\u5b9a\u7740\u3002TGC TEEN\u3067\u30c7\u30d3\u30e5\u30fc\u3002SNS\u30aa\u30fc\u30ac\u30cb\u30c3\u30af\u62e1\u6563\u3002\u30d5\u30a1\u30fc\u30b9\u30c8\u30b7\u30fc\u30ba\u30f3\u516820\u8a71\u5c55\u958b\u3002",
     RGBColor(0x1E, 0x6F, 0xBF)),
    ("PHASE 2", "2026-27", "\u77ed\u7de8\u6620\u753b\u30fb\u30e1\u30c7\u30a3\u30a2\u9023\u643a",
     "\u6e0b\u8c37\u821e\u53f0\u306e\u77ed\u7de8\u6620\u753b\u3002\u5730\u4e0a\u6ce2\u9023\u52d5\u4f01\u753b\u3002\u30df\u30b9\u30b3\u30f3\u76f8\u4e92\u9001\u5ba2\u3002\u30d5\u30a1\u30f3\u30af\u30e9\u30d6\u8a2d\u7acb\u3002",
     GOLD),
    ("PHASE 3", "2027-28", "IP\u78ba\u7acb\u30fbNo.1",
     "\u5730\u4e0a\u6ce2\u30ec\u30ae\u30e5\u30e9\u30fc\u3002\u5927\u624b\u30b5\u30d6\u30b9\u30afPF\u72ec\u5360\u914d\u4fe1\u3002\u30e1\u30f3\u30ba\u82b8\u80fd\u90e8\u9580\u30b7\u30a7\u30a21\u4f4d\u3002100\u5104\u5186IP\u4fa1\u5024\u3002",
     RED),
]
for i, (phase, year, title, desc, accent) in enumerate(phases):
    x = Inches(0.8 + i * 5.0)
    y = Inches(2.5)
    add_shape(slide, x, y, Inches(4.6), Inches(5.5), fill_color=CARD_BG)
    add_shape(slide, x, y, Inches(4.6), Inches(0.06), fill_color=accent)
    add_text(slide, x + Inches(0.3), y + Inches(0.3), Inches(4), Inches(0.4), phase, 14, accent, bold=True)
    add_text(slide, x + Inches(0.3), y + Inches(0.75), Inches(4), Inches(0.4), year, 16, GRAY)
    add_text(slide, x + Inches(0.3), y + Inches(1.3), Inches(4), Inches(0.6), title, 22, WHITE, bold=True)
    add_text(slide, x + Inches(0.3), y + Inches(2.2), Inches(4), Inches(3.0), desc, 16, LIGHT_GRAY)

# ─── SLIDE 9: KPI ───
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_shape(slide, Inches(0), Inches(0), Inches(0.12), H, fill_color=GOLD)
add_text(slide, Inches(0.8), Inches(0.5), Inches(5), Inches(0.6), "KPI", 14, GOLD, bold=True)

kpis = [
    ("\u00a5100\u5104", "IP VALUE TARGET", "2028\u5e74"),
    ("100\u4e07", "YouTube\u767b\u9332\u8005", ""),
    ("No.1", "\u30e1\u30f3\u30ba\u82b8\u80fd\u90e8\u9580", ""),
    ("300\u4e07/\u6708", "\u6708\u6b21\u5236\u4f5c\u4e88\u7b97", ""),
]
for i, (num, label, sub) in enumerate(kpis):
    x = Inches(0.8 + i * 3.7)
    y = Inches(1.5)
    add_shape(slide, x, y, Inches(3.4), Inches(2.5), fill_color=CARD_BG)
    add_text(slide, x + Inches(0.3), y + Inches(0.3), Inches(2.8), Inches(1.0), num, 36, GOLD, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, x + Inches(0.3), y + Inches(1.4), Inches(2.8), Inches(0.5), label, 16, WHITE, align=PP_ALIGN.CENTER)
    if sub:
        add_text(slide, x + Inches(0.3), y + Inches(1.9), Inches(2.8), Inches(0.4), sub, 13, GRAY, align=PP_ALIGN.CENTER)

add_text(slide, Inches(0.8), Inches(4.5), Inches(14), Inches(0.6),
         "Competitive Advantages", 20, WHITE, bold=True)

advantages = [
    ("\u2460  \u6e0b\u8c37\u3068\u3044\u3046\u5727\u5012\u7684\u30d6\u30e9\u30f3\u30c9", RED),
    ("\u2461  Z\u30fb\u03b1\u4e16\u4ee3\u5973\u6027\u306e\u6d88\u8cbb\u884c\u52d5\uff08\u63a8\u3057\u8ab2\u91d1\u6587\u5316\uff09", GOLD),
    ("\u2462  HJ\u306e\u72ec\u81ea\u30cd\u30c3\u30c8\u30ef\u30fc\u30af\uff08\u30df\u30b9\u30b3\u30f3 \u00d7 TGC TEEN \u00d7 \u5730\u4e0a\u6ce2\uff09", RGBColor(0x1E, 0x6F, 0xBF)),
]
for i, (text, accent) in enumerate(advantages):
    y = Inches(5.3 + i * 1.1)
    add_shape(slide, Inches(0.8), y, Inches(14.4), Inches(0.9), fill_color=CARD_BG)
    add_shape(slide, Inches(0.8), y, Inches(0.08), Inches(0.9), fill_color=accent)
    add_text(slide, Inches(1.2), y + Inches(0.15), Inches(13.8), Inches(0.6), text, 18, LIGHT_GRAY)

# ─── SLIDE 10: CLOSING ───
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_shape(slide, Inches(0), Inches(0), W, Inches(0.08), fill_color=RED)
add_multi_text(slide, Inches(1.5), Inches(2.5), Inches(13), Inches(4), [
    ("\u6e0b\u4e00\uff08\u30b7\u30d6\u30ef\u30f3\uff09", 52, WHITE, True, PP_ALIGN.CENTER),
    ("", 12, WHITE, False, PP_ALIGN.CENTER),
    ("\u65e5\u672c\u6700\u5927\u306e\u30e1\u30f3\u30baIP\u5275\u51fa\u88c5\u7f6e\u3078\u3002", 26, GRAY, False, PP_ALIGN.CENTER),
    ("", 20, WHITE, False, PP_ALIGN.CENTER),
    ("\u30c6\u30fc\u30de\u66f2:  O.K.thanks \u65b0\u66f2   |   \u30b9\u30dd\u30f3\u30b5\u30fc:  O.K.thanks", 18, GOLD, False, PP_ALIGN.CENTER),
])
add_shape(slide, Inches(5), Inches(6.2), Inches(6), Inches(0.04), fill_color=RED)
add_text(slide, Inches(1.5), Inches(7.5), Inches(13), Inches(0.6),
         "CONFIDENTIAL  \u00a9  HJ Inc. All Rights Reserved.", 14, GRAY, align=PP_ALIGN.CENTER)

# Save
out = os.path.join(os.path.dirname(os.path.abspath(__file__)), "\u6e0b\u4e00\u30d7\u30ed\u30b8\u30a7\u30af\u30c8\u6982\u8981.pptx")
prs.save(out)
print(f"Saved: {out}")
